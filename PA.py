import collections 
import collections.abc
from pptx import Presentation

# Crearea unei prezentări noi
presentation = Presentation()

# Slide 1 - Titlu
slide_layout = presentation.slide_layouts[0]
slide = presentation.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Prezentare Proiect Weather"
subtitle.text = "Principii SOLID și Design Pattern-uri"

# Slide 2 - Descrierea proiectului
slide_layout = presentation.slide_layouts[1]
slide = presentation.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Proiectul Weather"
content.text = "Proiectul Weather este o aplicație de prognoză meteo care permite utilizatorilor să obțină informații despre condițiile meteo curente și previziunile pentru diferite locații. Aplicația utilizează date în timp real și oferă o interfață intuitivă pentru utilizatori."

# Slide 3 - Principiul SOLID: Single Responsibility Principle (SRP)
slide_layout = presentation.slide_layouts[1]
slide = presentation.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Principiul SOLID: SRP"
content.text = "Single Responsibility Principle (SRP) este un principiu de proiectare care sugerează ca o clasă ar trebui să aibă o singură responsabilitate și să fie responsabilă pentru un singur aspect al sistemului. În proiectul Weather, SRP a fost aplicat prin separarea funcționalităților legate de procesarea datelor meteo, afișarea interfeței utilizatorului și gestionarea cererilor API în clase separate."

# Slide 4 - Principiul SOLID: Open-Closed Principle (OCP)
slide_layout = presentation.slide_layouts[1]
slide = presentation.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Principiul SOLID: OCP"
content.text = "Open-Closed Principle (OCP) este un principiu de proiectare care sugerează că o clasă trebuie să fie deschisă pentru extensie, dar închisă pentru modificare. În proiectul Weather, OCP a fost aplicat prin utilizarea de interfețe și clase abstracte pentru a permite extinderea funcționalității fără a modifica clasele existente. De exemplu, adăugarea unui nou furnizor de date meteo nu necesită modificarea codului existent, ci doar implementarea unei noi clase care respectă interfața comună."

# Slide 5 - Principiul SOLID: Liskov Substitution Principle (LSP)
slide_layout = presentation.slide_layouts[1]
slide = presentation.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Principiul SOLID: LSP"
content.text = "Liskov Substitution Principle (LSP) este un principiu de proiectare care afirmă că obiectele de tipul unei clase de bază pot fi înlocuite cu obiecte de tipul unei subclase, fără a afecta corectitudinea programului. În proiectul Weather, LSP a fost respectat prin utilizarea de ierarhii de clase și moștenire, astfel încât obiectele specifice, cum ar fi un furnizor de date meteo sau un procesor de date, pot fi folosite în locul obiectelor generice."

# Slide 6 - Principiul SOLID: Interface Segregation Principle (ISP)
slide_layout = presentation.slide_layouts[1]
slide = presentation.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Principiul SOLID: ISP"
content.text = "Interface Segregation Principle (ISP) este un principiu de proiectare care afirmă că clienții nu trebuie să fie forțați să depindă de interfețe pe care nu le utilizează. În proiectul Weather, ISP a fost aplicat prin crearea de interfețe mici și specializate care sunt implementate doar de clasele care le folosesc efectiv. De exemplu, interfața pentru un procesor de date meteo conține doar metodele relevante pentru procesarea datelor, evitând astfel dependențe inutile."

# Slide 7 - Design Pattern: Singleton
slide_layout = presentation.slide_layouts[1]
slide = presentation.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Design Pattern: Singleton"
content.text = "Singleton este un design pattern creational care permite crearea unei singure instanțe a unei clase și furnizează un punct global de acces la acea instanță. În proiectul Weather, Singleton a fost utilizat pentru a crea o instanță unică a unui manager de configurații, care gestionează setările aplicației și oferă acces la acestea în întregul sistem."

# Slide 8 - Design Pattern: Strategy
slide_layout = presentation.slide_layouts[1]
slide = presentation.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Design Pattern: Strategy"
content.text = "Strategy este un design pattern behavioral care permite definirea unei familii de algoritmi, încapsularea fiecărui algoritm și făcându-i interschimbabili. În proiectul Weather, Strategy a fost utilizat pentru a implementa diferite strategii de procesare a datelor meteo, cum ar fi procesarea în timp real, procesarea istorică sau procesarea în mod offline. Aceste strategii pot fi schimbate în timpul execuției fără a afecta restul sistemului."

# Slide 9 - Design Pattern: Observer
slide_layout = presentation.slide_layouts[1]
slide = presentation.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Design Pattern: Observer"
content.text = "Observer este un design pattern behavioral care permite notificarea automată a unor obiecte dependente de schimbările de stare ale unui obiect observabil. În proiectul Weather, Observer a fost utilizat pentru a implementa funcționalitatea de notificare a utilizatorilor despre actualizările meteo. Utilizatorii pot fi înregistrați ca observatori și vor fi notificați automat atunci când se produc schimbări în datele meteo."

# Slide 10 - Încheiere
slide_layout = presentation.slide_layouts[1]
slide = presentation.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Încheiere"
content.text = "Prin respectarea principiilor SOLID și utilizarea design pattern-urilor potrivite, proiectul Weather a obținut o arhitectură robustă, ușor de întreținut și extensibilă. Aceste principii și pattern-uri oferă flexibilitate, modularitate și separare a responsabilităților, ceea ce facilitează dezvoltarea și evoluția aplicației pe termen lung."

# Salvarea prezentării
presentation.save("prezentare_weather.pptx")
